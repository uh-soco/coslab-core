import sqlite3
from nltk.corpus import wordnet as wn
from nltk.corpus import wordnet_ic as wn_ic
import itertools
from wordcloud import WordCloud
import numpy as np
import collections
from pptx import Presentation
from pptx.util import Inches

# prs = Presentation()


def dict_factory(cursor, row):
    d = {}
    for idx, col in enumerate(cursor.description):
        d[col[0]] = row[idx]
    return d


##cache results for efficiency
try:
    all_similarities
except:
    all_similarities = {}


def evaluate_similarity(synset1, synset2):

    # synsets1 = wn.synsets( word1 )
    # synsets2 = wn.synsets( word2 )

    best = 0

    for s1, s2 in itertools.product(synsets1, synsets2):

        score = -1

        if not (s1, s2) in all_similarities:
            try:
                _score = s1.wup_similarity(s2) or -1
                all_similarities[(s1, s2)] = _score
                all_similarities[(s2, s1)] = _score
            except:
                pass

        score = all_similarities[(s2, s1)]

        if score > best:
            best = score

    return best


##########################################################################################################################
# MAKE DATA READY FOR ANALYSIS

conn = sqlite3.connect("results.db")
conn.row_factory = dict_factory
db = conn.cursor()

image_ids = collections.defaultdict(list)

for result in db.execute("select * from results"):
    image_ids[result["image"]].append(result["id"])
print(len(image_ids.keys()))

labels_per_service = collections.defaultdict(list)
labels_per_image_service = collections.defaultdict(
    lambda: collections.defaultdict(list)
)

for image, ids in image_ids.items():
    print(", ".join(map(str, ids)))
    for result in db.execute(
        "SELECT (SELECT service FROM results where id==label) as service, * FROM results WHERE id in (%s)"
        % (", ".join(map(str, ids)))
    ):
        labels_per_service[result["service"]].append(result["label"])
        labels_per_image_service[image][result["service"]].append(result["label"])

#############################################################################################################################
# WORD CLOUDS TO VISUALIZE DATA


for service, labels in labels_per_service.items():
    wordcloud = WordCloud(
        max_font_size=50, max_words=100, background_color="white"
    ).generate(" ".join(labels))
    wordcloud.to_file("fb_" + service + ".png")

##############################################################################################################################3
# ACCURACY METRICS

services = labels_per_service.keys()

similarity = collections.defaultdict(int)
max_change = collections.defaultdict(int)

for image, d in labels_per_image_service.items():

    for service1 in services:
        for service2 in services:
            ## compare services 1 and 2

            d1 = set(d[service1])
            d2 = set(d[service2])

            ## similarity
            sim = len(d1.intersection(d2))

            max_change[service1, service2] += min(len(d1), len(d2))
            similarity[service1, service2] += sim

for service1 in services:
    for service2 in services:

        accuracy = similarity[service1, service2] / max_change[service1, service2]

        print(service1, service2, accuracy)

# slide = prs.slides.add_slide(prs.slide_layouts[5])
# x, y, cx, cy = Inches(2), Inches(2), Inches(4), Inches(1.5)
# shape = slide.shapes.add_table(5, 5, x, y, cx, cy)

# for i, service1 in enumerate(services):
#     for j, service2 in enumerate(services):

#         accuracy = similarity[service1, service2] / max_change[service1, service2]

#         shape.table.cell(i + 1, 0).text = service1
#         shape.table.cell(0, j + 1).text = service2

#         shape.table.cell(i + 1, j + 1).text = "{:2.2f}".format(accuracy * 100)

#################################################################################################################
# EVALUATE SYNONYMS


brown_ic = wn_ic.ic("ic-brown.dat")

synsets_for_service = collections.defaultdict(list)
synsets_for_image_service = collections.defaultdict(
    lambda: collections.defaultdict(list)
)

for image, ids in image_ids.items():
    print(", ".join(map(str, ids)))
    for result in db.execute(
        "SELECT (SELECT service FROM results where id==label) as service, * FROM results WHERE id in (%s)"
        % (", ".join(map(str, ids)))
    ):

        synsets_for_service[result["service"]] += wn.synsets(result["label"])
        synsets_for_image_service[image][result["service"]] += wn.synsets(
            result["label"]
        )


similarity = collections.defaultdict(list)

for image in synsets_for_image_service.keys():

    for service1, service2 in itertools.permutations(services, 2):
        synsets1 = synsets_for_image_service[image][service1]
        synsets2 = synsets_for_image_service[image][service2]

        score = evaluate_similarity(synsets1, synsets2)

        similarity[service1, service2].append(score)

for services, paths in similarity.items():

    best = max(paths)
    mean = np.mean(paths)
    median = np.median(paths)

    print(services[0], services[1], "median", median)

for services, paths in similarity.items():

    best = max(paths)
    mean = np.mean(paths)
    median = np.median(paths)

    print(services[0], services[1], "mean", mean)

# slide = prs.slides.add_slide(prs.slide_layouts[5])
# x, y, cx, cy = Inches(2), Inches(2), Inches(4), Inches(1.5)
# shape = slide.shapes.add_table(5, 5, x, y, cx, cy)

# services = ["aws", "azure", "googlecloud"]

# for i, service1 in enumerate(services):
#     for j, service2 in enumerate(services):

#         if service1 != service2:

#             paths = similarity[service1, service2]

#             best = max(paths)
#             mean = np.mean(paths)
#             median = np.median(paths)

#             shape.table.cell(i + 1, 0).text = service1
#             shape.table.cell(0, j + 1).text = service2

#             shape.table.cell(i + 1, j + 1).text = "{:2.2f}".format(mean * 100)

# prs.save("summary.pptx")
